#!/usr/bin/env python3
"""
Build scripts/pybamm_comparison_slides.pptx  — 3-slide academic deck.
Theme: navy (#1A2E5A) + amber (#F5A623), white background, clean academic layout.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent.parent
IMG_BIAS = ROOT / "scripts" / "pybamm_vs_opencathode_bias2pct.png"
IMG_BASE = ROOT / "scripts" / "pybamm_vs_opencathode_all.png"
OUT_PATH = ROOT / "scripts" / "pybamm_comparison_slides.pptx"

# ── palette ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x2E, 0x5A)
AMBER   = RGBColor(0xF5, 0xA6, 0x23)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
COAL    = RGBColor(0x22, 0x22, 0x22)
GRAY    = RGBColor(0x88, 0x88, 0x88)
NAVY_LT = RGBColor(0xE8, 0xEE, 0xF6)   # table alt-row tint
AMBER_T = RGBColor(0xFF, 0xF0, 0xCC)   # amber table header tint
GREEN   = RGBColor(0x1E, 0x7A, 0x3E)
RED     = RGBColor(0xB5, 0x1F, 0x1F)


def _rgb_fill(shape, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    sp_pr = shape.fill
    sp_pr.solid()
    sp_pr.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, fill: RGBColor) -> object:
    from pptx.util import Emu
    rect = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    rect.line.fill.background()
    _rgb_fill(rect, fill)
    return rect


def _txbox(slide, x, y, w, h) -> object:
    return slide.shapes.add_textbox(x, y, w, h)


def _tf(shape, text="", size=11, bold=False, color=COAL,
        align=PP_ALIGN.LEFT, italic=False, space_before=0) -> object:
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def _add_para(tf, text="", size=11, bold=False, color=COAL,
              align=PP_ALIGN.LEFT, italic=False, space_before=4) -> None:
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _title_bar(slide, title: str, sw: Emu, bar_h=Inches(1.15)) -> None:
    """Navy bar spanning full width, white bold title text."""
    bar = _add_rect(slide, Emu(0), Emu(0), sw, bar_h, NAVY)
    tb = _txbox(slide, Inches(0.30), Inches(0.18), sw - Inches(0.6), Inches(0.80))
    tf = _tf(tb, title, size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    tf.word_wrap = True


def _footer(slide, sw: Emu, sh: Emu, text: str) -> None:
    tb = _txbox(slide, Inches(0.30), sh - Inches(0.35), sw - Inches(0.60), Inches(0.30))
    _tf(tb, text, size=7, color=GRAY, italic=True)


def _amber_rule(slide, x, y, w, h=Inches(0.04)) -> None:
    rule = _add_rect(slide, x, y, w, h, AMBER)


# ── table helper ─────────────────────────────────────────────────────────────

def _set_cell(cell, text, size=10, bold=False, color=COAL,
              fill: RGBColor | None = None,
              align=PP_ALIGN.LEFT) -> None:
    if fill:
        _rgb_fill(cell, fill)
    cell.margin_left  = Pt(6)
    cell.margin_right = Pt(6)
    cell.margin_top   = Pt(4)
    cell.margin_bottom = Pt(4)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _build_deck() -> None:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)
    SW = prs.slide_width
    SH = prs.slide_height

    blank = prs.slide_layouts[6]   # completely blank

    # =========================================================================
    # SLIDE 1 — "Different Tools for Different Problems"
    # =========================================================================
    s1 = prs.slides.add_slide(blank)
    _title_bar(s1, "PyBaMM vs OpenCATHODE: Different Tools for Different Problems", SW)
    _amber_rule(s1, Emu(0), Inches(1.15), SW, Inches(0.045))

    # Subtitle
    sb = _txbox(s1, Inches(0.30), Inches(1.28), SW - Inches(0.60), Inches(0.40))
    _tf(sb, "The fair test is the real BMS task: online SOC estimation from noisy field data.",
        size=12, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Left column: PyBaMM ──────────────────────────────────────────────────
    lx, ly, lw, lh = Inches(0.30), Inches(1.80), Inches(5.90), Inches(4.80)
    lb = _add_rect(s1, lx, ly, lw, lh, NAVY_LT)

    hdr_l = _txbox(s1, lx + Inches(0.15), ly + Inches(0.15), lw - Inches(0.30), Inches(0.45))
    _tf(hdr_l, "PyBaMM  —  Forward Simulator", size=14, bold=True, color=NAVY)

    body_l = _txbox(s1, lx + Inches(0.15), ly + Inches(0.65), lw - Inches(0.30), Inches(3.90))
    tf_l = _tf(body_l,
               "Input:  known cell parameters  +  applied current I(t)",
               size=11, color=COAL)
    _add_para(tf_l, "Output:  predicted voltage V(t)  and  open-loop SOC(t)", size=11, color=COAL, space_before=6)
    _add_para(tf_l, "", size=6, color=COAL)
    _add_para(tf_l, "No voltage feedback.  Init SOC error propagates indefinitely.", size=11, bold=True, color=RED, space_before=4)
    _add_para(tf_l, "", size=6)
    _add_para(tf_l, "Best suited for:", size=11, bold=True, color=NAVY, space_before=4)
    _add_para(tf_l, "  · Controlled lab cells with exact known parameters", size=11, color=COAL, space_before=2)
    _add_para(tf_l, "  · Electrode design and parameter sensitivity studies", size=11, color=COAL, space_before=2)
    _add_para(tf_l, "  · DFN physics research — not field deployment", size=11, color=COAL, space_before=2)

    # ── Right column: OpenCATHODE ────────────────────────────────────────────
    rx, ry, rw, rh = Inches(6.90), Inches(1.80), Inches(6.10), Inches(4.80)
    rb = _add_rect(s1, rx, ry, rw, rh, RGBColor(0xF0, 0xF5, 0xFF))

    hdr_r = _txbox(s1, rx + Inches(0.15), ry + Inches(0.15), rw - Inches(0.30), Inches(0.45))
    _tf(hdr_r, "OpenCATHODE  —  Inverse Estimator", size=14, bold=True, color=NAVY)

    body_r = _txbox(s1, rx + Inches(0.15), ry + Inches(0.65), rw - Inches(0.30), Inches(3.90))
    tf_r = _tf(body_r,
               "Input:  noisy field current  +  measured voltage V_meas(t)",
               size=11, color=COAL)
    _add_para(tf_r, "Output:  recovered hidden SOC(t) — closed-loop Kalman update", size=11, color=COAL, space_before=6)
    _add_para(tf_r, "", size=6)
    _add_para(tf_r, "V_meas feedback corrects SOC drift at every timestep.", size=11, bold=True, color=GREEN, space_before=4)
    _add_para(tf_r, "", size=6)
    _add_para(tf_r, "Best suited for:", size=11, bold=True, color=NAVY, space_before=4)
    _add_para(tf_r, "  · Unknown initial SOC, biased sensors, chemistry variation", size=11, color=COAL, space_before=2)
    _add_para(tf_r, "  · Fleet vehicles with heterogeneous cell ageing", size=11, color=COAL, space_before=2)
    _add_para(tf_r, "  · Real-time 1 Hz embedded BMS  (54 µs/cell, 36-cell pack)", size=11, color=COAL, space_before=2)

    # Amber vertical divider
    _amber_rule(s1, Inches(6.63), Inches(1.80), Inches(0.04), Inches(4.80))

    _footer(s1, SW, SH,
            "Sharma 2026 · OpenCATHODE Stack · Validated: 634 450 datapoints, Quartz WLTP NMC811 "
            "· MAE 18.6 mV, R²=0.981")

    # =========================================================================
    # SLIDE 2 — "Three Structural Advantages (with proof)"
    # =========================================================================
    s2 = prs.slides.add_slide(blank)
    _title_bar(s2, "Three Structural Advantages (with proof)", SW)
    _amber_rule(s2, Emu(0), Inches(1.15), SW, Inches(0.045))

    # ── Comparison table (left 6.7") ─────────────────────────────────────────
    TBL_X = Inches(0.30)
    TBL_Y = Inches(1.28)
    TBL_W = Inches(6.70)
    TBL_H = Inches(3.00)
    COL_W = [Inches(2.55), Inches(2.10), Inches(2.05)]

    tbl = s2.shapes.add_table(4, 3, TBL_X, TBL_Y, TBL_W, TBL_H).table
    tbl.columns[0].width = COL_W[0]
    tbl.columns[1].width = COL_W[1]
    tbl.columns[2].width = COL_W[2]

    # Header row
    _set_cell(tbl.cell(0, 0), "Dimension",             size=10, bold=True, color=WHITE, fill=NAVY, align=PP_ALIGN.CENTER)
    _set_cell(tbl.cell(0, 1), "OpenCATHODE EKF",       size=10, bold=True, color=WHITE, fill=NAVY, align=PP_ALIGN.CENTER)
    _set_cell(tbl.cell(0, 2), "PyBaMM  (Chen2020)",    size=10, bold=True, color=WHITE, fill=NAVY, align=PP_ALIGN.CENTER)

    # Data rows
    rows = [
        ("Voltage prediction\nV MAE (BMW / VED / Quartz)",
         "7 mV / 27 mV / 73 mV",
         "38.7 mV / 182 mV / 472 mV"),
        ("Current-bias immunity\n(Δ SOC error at +2% bias)",
         "≤ 0.1 pp  (bias absorbed\nby V_meas feedback)",
         "Drifts monotonically\n−0.7 pp @ Quartz, 72%\ndirectional steps"),
        ("Applicability\n(Deng BAIC CC charging)",
         "Runs  ✓",
         "Solver abort:\nMax-voltage event\nviolated at t=0"),
    ]
    row_fills = [WHITE, NAVY_LT]
    val_colors_ekf = [GREEN, GREEN, GREEN]
    val_colors_pb  = [RED,   RED,   RED]
    for i, (dim, ekf_val, pb_val) in enumerate(rows):
        fill = row_fills[i % 2]
        _set_cell(tbl.cell(i+1, 0), dim,     size=9,  bold=False, color=COAL,                fill=fill)
        _set_cell(tbl.cell(i+1, 1), ekf_val, size=9,  bold=True,  color=val_colors_ekf[i], fill=fill)
        _set_cell(tbl.cell(i+1, 2), pb_val,  size=9,  bold=True,  color=val_colors_pb[i],  fill=fill)

    # ── Three notes below table ───────────────────────────────────────────────
    notes_box = _txbox(s2, TBL_X, TBL_Y + TBL_H + Inches(0.15),
                       TBL_W, Inches(2.65))
    tf_n = _tf(notes_box, "① Voltage:  5–6× better  —  PyBaMM requires exact per-cell DFN params "
               "(unavailable in the field); EKF adapts via V_meas online.",
               size=9, color=COAL, italic=False)
    _add_para(tf_n, "", size=4)
    _add_para(tf_n, "② Bias drift:  open-loop SOC is unbounded under sensor bias; "
              "closed-loop Kalman update cancels it at every step.",
              size=9, color=COAL, space_before=2)
    _add_para(tf_n, "", size=4)
    _add_para(tf_n, "③ Solver abort:  Chen2020 (LG M50 cell, lab conditions) cannot simulate "
              "a real CC-charging profile at SOC=27%  —  "
              "field applicability is fundamentally limited.",
              size=9, color=COAL, space_before=2)

    # ── Image: SOC-error panel (right side) ───────────────────────────────────
    IMG_X = Inches(7.15)
    IMG_Y = Inches(1.28)
    IMG_W = Inches(5.85)          # AR≈1.96 → height = 5.85/1.96 ≈ 2.98"
    IMG_H = Inches(2.99)
    s2.shapes.add_picture(str(IMG_BIAS), IMG_X, IMG_Y, IMG_W, IMG_H)

    # Caption under image
    cap2 = _txbox(s2, IMG_X, IMG_Y + IMG_H + Inches(0.06), IMG_W, Inches(0.35))
    _tf(cap2,
        "SOC error vs time (bias=2%).  Row 2 shows PyBaMM drift (red, monotone ↓) "
        "vs EKF bounded (blue). Shaded = |PB| > |EKF|.",
        size=7.5, color=GRAY, italic=True)

    _footer(s2, SW, SH,
            "All numbers from scripts/compare_pybamm_all.py  —  real field data, single-trip "
            "mode, no fleet calibration  · PyBaMM 26.6.2.0 · Chen2020 param set")

    # =========================================================================
    # SLIDE 3 — "What I Openly Concede (the honest core)"
    # =========================================================================
    s3 = prs.slides.add_slide(blank)
    _title_bar(s3, "What I Openly Concede (the honest core)", SW)
    _amber_rule(s3, Emu(0), Inches(1.15), SW, Inches(0.045))

    # ── Bullet content (left 55%) ─────────────────────────────────────────────
    BODY_X = Inches(0.30)
    BODY_Y = Inches(1.30)
    BODY_W = Inches(6.65)
    BODY_H = Inches(5.75)

    body3 = _txbox(s3, BODY_X, BODY_Y, BODY_W, BODY_H)
    tf3 = _tf(body3,
              "① Clean-start accuracy: PyBaMM wins",
              size=13, bold=True, color=NAVY)
    _add_para(tf3,
              "With correct initial SOC and no sensor bias, PyBaMM's open-loop "
              "Coulomb counting achieves 0.1–1.3% SOC RMSE.  My single-trip EKF "
              "shows 11–18% RMSE — not because it's wrong, but because the OCV "
              "model mismatch moves the Kalman state away from the true SOC.",
              size=10.5, color=COAL, space_before=4)
    _add_para(tf3, "", size=6)

    _add_para(tf3,
              "② My EKF's value is robustness — not clean-case accuracy",
              size=13, bold=True, color=NAVY, space_before=6)
    _add_para(tf3,
              "The EKF is designed for unknown initial SOC, biased sensors, and "
              "cell-to-cell variation over a vehicle's lifetime.  In that scenario "
              "PyBaMM's error grows without bound; the EKF's is bounded by "
              "V_meas.  They don't compete — they solve different problems.",
              size=10.5, color=COAL, space_before=4)
    _add_para(tf3, "", size=6)

    _add_para(tf3,
              "③ Fleet SOC accuracy is OCV-model-limited (sourced numbers)",
              size=13, bold=True, color=NAVY, space_before=6)
    _add_para(tf3,
              "Mode B (free-running EKF, +20% init offset, PCHIP cal, held-out 90%): "
              "Deng BAIC 11.9% · BMW i3 20.8% · VED 25.5% SOC RMSE.  "
              "Spread reflects OCV-model chemistry mismatch, not filter instability — "
              "the EKF converges and stays bounded; accuracy is gated by OCV table quality.",
              size=10.5, color=COAL, space_before=4)
    _add_para(tf3, "", size=6)

    # Amber rule above bottom note
    _amber_rule(s3, BODY_X, BODY_Y + BODY_H - Inches(0.75), BODY_W, Inches(0.035))
    bottom_note = _txbox(s3, BODY_X, BODY_Y + BODY_H - Inches(0.65), BODY_W, Inches(0.55))
    _tf(bottom_note,
        "Takeaway: use PyBaMM for lab physics research; use OpenCATHODE EKF for "
        "real-time field BMS where sensor noise, unknown SOC_0, and cell variation dominate.",
        size=10, bold=True, italic=True, color=AMBER,
        align=PP_ALIGN.LEFT)

    # ── Image: baseline comparison (right side) ───────────────────────────────
    IM3_X = Inches(7.20)
    IM3_Y = Inches(1.30)
    IM3_W = Inches(5.80)
    IM3_H = Inches(2.96)
    s3.shapes.add_picture(str(IMG_BASE), IM3_X, IM3_Y, IM3_W, IM3_H)

    cap3 = _txbox(s3, IM3_X, IM3_Y + IM3_H + Inches(0.06), IM3_W, Inches(0.28))
    _tf(cap3,
        "Baseline (no init offset, no bias).  PyBaMM SOC (row 1) tracks truth "
        "at 0.1–1.3% RMSE; EKF drifts due to OCV model mismatch without fleet cal.",
        size=7.5, color=GRAY, italic=True)

    # Key stat box (amber) - right side below image
    # stat_y = IM3_Y + IM3_H + 0.40 so bottom = 4.66+2.45=7.11" < footer at 7.15"
    stat_x = IM3_X
    stat_y = IM3_Y + IM3_H + Inches(0.40)
    stat_w = IM3_W
    stat_h = Inches(2.45)
    stat_bg = _add_rect(s3, stat_x, stat_y, stat_w, stat_h, AMBER_T)
    stat_box = _txbox(s3, stat_x + Inches(0.15), stat_y + Inches(0.12),
                      stat_w - Inches(0.30), stat_h - Inches(0.20))
    tf_stat = _tf(stat_box, "Fleet-validated performance",
                  size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _add_para(tf_stat, "", size=4)
    stats = [
        ("Voltage MAE (Quartz WLTP, 36 cells)", "18.6 mV  ✓  < 20 mV target"),
        ("Fleet SOC RMSE (Deng / BMW / VED)",   "11.9% / 20.8% / 25.5%  — OCV-limited"),
        ("Current-bias Δ at +2%",               "≤ 0.1 pp  (EKF)  vs  0.7 pp  (PyBaMM)"),
        ("Step latency  (54 µs/cell)",           "real-time 1 Hz  ✓"),
        ("Bias-immune via",                      "V_meas closed-loop Kalman update"),
    ]
    for label, val in stats:
        _add_para(tf_stat, f"  {label}:", size=9, bold=True, color=NAVY, space_before=4)
        _add_para(tf_stat, f"    {val}", size=9.5, bold=False, color=GREEN, space_before=1)

    _footer(s3, SW, SH,
            "Full validation: data/validate_quartz.py  ·  Fleet sweep: data/validate_generic.py  "
            "·  Audit: scripts/audit_independent.py  ·  Sharma 2026")

    # ── save ─────────────────────────────────────────────────────────────────
    prs.save(str(OUT_PATH))
    print(f"[OK] Saved → {OUT_PATH}")
    return OUT_PATH


def _verify_layout(path: Path) -> None:
    """Load saved PPTX and verify no shape exceeds slide bounds or invades footer zone."""
    prs = Presentation(str(path))
    SW = prs.slide_width
    SH = prs.slide_height
    FOOTER_Y = SH - Inches(0.35)   # footer top; content must end above this
    ok = True
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            r = shape.left + shape.width
            b = shape.top  + shape.height
            overflow_r = r > SW + Inches(0.05)
            overflow_b = b > SH + Inches(0.05)
            # footer shapes sit near the bottom with small height — skip them
            is_footer = (shape.top > SH - Inches(0.50)
                         and shape.height < Inches(0.40))
            invades_footer = (not is_footer and b > FOOTER_Y + Inches(0.01))
            if overflow_r or overflow_b:
                print(f"  [WARN] Slide {i} '{shape.name}': "
                      f"right={r/914400:.2f}\" bottom={b/914400:.2f}\"  "
                      f"(slide {SW/914400:.2f}\" × {SH/914400:.2f}\")")
                ok = False
            elif invades_footer:
                print(f"  [WARN] Slide {i} '{shape.name}': bottom={b/914400:.3f}\" "
                      f"invades footer zone (>{FOOTER_Y/914400:.3f}\")")
                ok = False
    if ok:
        print("[VERIFY] No shape overflows or footer collisions  ✓")
    else:
        print("[VERIFY] Issue(s) detected — check slide layout.")


if __name__ == "__main__":
    for img in [IMG_BIAS, IMG_BASE]:
        if not img.exists():
            print(f"[ERROR] Missing image: {img}")
            sys.exit(1)
    out = _build_deck()
    _verify_layout(Path(out) if isinstance(out, str) else OUT_PATH)
